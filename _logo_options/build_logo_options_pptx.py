"""Build the logo-options A/B comparison PPTX.

Renders each candidate SVG at multiple sizes alongside the current mark
so the team can compare clarity, simplicity, and favicon-scale
performance. Output: ../../corp docs/NeedleGPS_Logo_Options.pptx

Run: /Users/lkini/Projects/NeedleGPS/.venv/bin/python build_logo_options_pptx.py
"""

from __future__ import annotations

import io
import subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


HERE = Path(__file__).resolve().parent
WEBSITE = HERE.parent
LOGO_DIR = HERE
CURRENT_MARK = WEBSITE / "assets" / "logos" / "needlegps-mark.png"

# Brand palette
BG_LIGHT = RGBColor(0xF6, 0xF8, 0xFA)
BG_CARD  = RGBColor(0xFF, 0xFF, 0xFF)
INK      = RGBColor(0x0E, 0x20, 0x30)
INK_SOFT = RGBColor(0x3B, 0x4A, 0x60)
INK_MUTE = RGBColor(0x6A, 0x76, 0x89)
INK_LINE = RGBColor(0xE2, 0xE6, 0xEC)
C_PRIM   = RGBColor(0x1F, 0xB6, 0xC9)

DISPLAY = "Manrope"
MONO    = "JetBrains Mono"

SLIDE_W_IN, SLIDE_H_IN = 16.0, 9.0


# rasterize SVG to PNG via macOS sips fallback to qlmanage
def svg_to_png(svg_path: Path, width: int) -> bytes:
    """Render an SVG to a PNG of the given pixel width. Uses macOS rsvg if
    available; falls back to sips which can also rasterize SVG on recent
    macOS versions."""
    tmp_png = svg_path.with_suffix(f".tmp{width}.png")
    # Try rsvg-convert first (best SVG renderer if installed via brew)
    try:
        subprocess.run(
            ["rsvg-convert", "-w", str(width), "-o", str(tmp_png), str(svg_path)],
            check=True, capture_output=True,
        )
    except (FileNotFoundError, subprocess.CalledProcessError):
        # Fall back to sips (built-in on macOS but won't always handle SVG)
        try:
            subprocess.run(
                ["sips", "-s", "format", "png", "-Z", str(width),
                 str(svg_path), "--out", str(tmp_png)],
                check=True, capture_output=True,
            )
        except Exception:
            # Last resort: use Pillow with cairosvg if available
            import cairosvg  # type: ignore
            cairosvg.svg2png(url=str(svg_path), write_to=str(tmp_png),
                             output_width=width)
    data = tmp_png.read_bytes()
    tmp_png.unlink(missing_ok=True)
    return data


def add_text(slide, x, y, w, h, text, *, font=DISPLAY, size=14, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def slide_chrome(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.fill.background()
    return slide


def slide_cover(prs):
    s = slide_chrome(prs)
    add_text(s, Inches(0.9), Inches(0.9), Inches(8), Inches(0.5),
             "NEEDLEGPS  ·  LOGO OPTIONS  ·  A/B",
             font=MONO, size=12, color=C_PRIM)
    add_text(s, Inches(0.9), Inches(1.3), Inches(14), Inches(2),
             "Six simpler marks for review.",
             font=DISPLAY, size=46, bold=True, color=INK)
    add_text(s, Inches(0.9), Inches(3.3), Inches(13), Inches(3.5),
             "Each option is rendered at three sizes (192 px, 64 px, 32 px)\n"
             "alongside the current mark so we can judge clarity at hero,\n"
             "nav, and favicon scales. Single-color monochrome by default\n"
             "to keep print + dark-mode + small-size legibility solid.",
             font=DISPLAY, size=18, color=INK_SOFT)
    add_text(s, Inches(0.9), Inches(7.8), Inches(14), Inches(0.4),
             "Voting: pick one + any tweaks. We swap in 5 minutes.",
             font=MONO, size=11, color=INK_MUTE)


# embed three sizes of a logo in a panel
def logo_panel(slide, x, y, w, h, label, svg_path, sub=None, *, is_png=False):
    # panel background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.06
    card.fill.solid(); card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = INK_LINE
    card.line.width = Pt(0.75)
    # label
    add_text(slide, x + Inches(0.25), y + Inches(0.2), w - Inches(0.5), Inches(0.4),
             label, font=MONO, size=11, color=C_PRIM)
    if sub:
        add_text(slide, x + Inches(0.25), y + Inches(0.55), w - Inches(0.5), Inches(0.4),
                 sub, font=DISPLAY, size=12, color=INK_SOFT)

    # 3 raster previews at increasing scales
    sizes = [(192, Inches(1.9)), (64, Inches(0.9)), (32, Inches(0.5))]
    cx = x + Inches(0.25)
    cy = y + h - Inches(2.2)
    if is_png:
        # current mark is already a PNG
        png_bytes = svg_path.read_bytes()
    else:
        png_bytes = None  # render below

    for idx, (px, render_w) in enumerate(sizes):
        if is_png:
            buf = png_bytes
        else:
            buf = svg_to_png(svg_path, px)
        bio = io.BytesIO(buf)
        slide.shapes.add_picture(bio, cx, cy, height=render_w, width=render_w)
        add_text(slide, cx, cy + render_w + Inches(0.05), render_w, Inches(0.3),
                 f"{px}×{px} px", font=MONO, size=9, color=INK_MUTE,
                 align=PP_ALIGN.CENTER)
        cx = cx + render_w + Inches(0.25)


def slide_compare(prs):
    """One slide showing all candidates + current side-by-side at one size."""
    s = slide_chrome(prs)
    add_text(s, Inches(0.9), Inches(0.5), Inches(8), Inches(0.4),
             "01 · ALL OPTIONS AT 96 × 96 PX",
             font=MONO, size=11, color=C_PRIM)
    add_text(s, Inches(0.9), Inches(0.9), Inches(14), Inches(1),
             "Side by side.",
             font=DISPLAY, size=34, bold=True, color=INK)

    entries = [
        ("CURRENT", "syringe + target (GitHub avatar)", CURRENT_MARK, True),
        ("A", "bare reticle", LOGO_DIR / "option-A-bare-reticle.svg", False),
        ("B", "crosshair", LOGO_DIR / "option-B-crosshair.svg", False),
        ("C", "needle arrow", LOGO_DIR / "option-C-needle-arrow.svg", False),
        ("D", "bubble level", LOGO_DIR / "option-D-bubble-level.svg", False),
        ("E", "trajectory", LOGO_DIR / "option-E-trajectory.svg", False),
        ("F", "N monogram", LOGO_DIR / "option-F-monogram.svg", False),
    ]
    cell_w = Inches(1.9); gap = Inches(0.18)
    x0 = Inches(0.9); y0 = Inches(2.8)
    img_size = Inches(1.4)
    for i, (label, sub, path, is_png) in enumerate(entries):
        x = x0 + (cell_w + gap) * i
        # tile background
        tile = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y0, cell_w, Inches(3.6))
        tile.adjustments[0] = 0.06
        tile.fill.solid(); tile.fill.fore_color.rgb = BG_CARD
        tile.line.color.rgb = INK_LINE; tile.line.width = Pt(0.75)
        # label
        add_text(s, x, y0 + Inches(0.25), cell_w, Inches(0.4),
                 label, font=MONO, size=11, color=C_PRIM,
                 align=PP_ALIGN.CENTER)
        # render the mark
        if is_png:
            png = path.read_bytes()
        else:
            png = svg_to_png(path, 256)
        bio = io.BytesIO(png)
        # center horizontally inside the tile
        s.shapes.add_picture(
            bio,
            x + (cell_w - img_size) / 2,
            y0 + Inches(0.7),
            width=img_size, height=img_size,
        )
        # subtitle
        add_text(s, x, y0 + Inches(2.4), cell_w, Inches(1),
                 sub, font=DISPLAY, size=11, color=INK_SOFT,
                 align=PP_ALIGN.CENTER)


def slide_for_option(prs, idx, label, name, svg, rationale):
    s = slide_chrome(prs)
    add_text(s, Inches(0.9), Inches(0.5), Inches(8), Inches(0.4),
             f"0{idx + 2} · OPTION {label}",
             font=MONO, size=11, color=C_PRIM)
    add_text(s, Inches(0.9), Inches(0.9), Inches(14), Inches(1.2),
             name, font=DISPLAY, size=42, bold=True, color=INK)
    add_text(s, Inches(0.9), Inches(2.4), Inches(8), Inches(2.5),
             rationale, font=DISPLAY, size=15, color=INK_SOFT)

    # three sizes on the right
    sizes = [(192, Inches(2.0)), (64, Inches(0.85)), (32, Inches(0.42))]
    cx = Inches(9.5); cy = Inches(2.4)
    for px, render_w in sizes:
        bio = io.BytesIO(svg_to_png(svg, px))
        s.shapes.add_picture(bio, cx, cy, height=render_w, width=render_w)
        add_text(s, cx, cy + render_w + Inches(0.05), render_w, Inches(0.3),
                 f"{px}px", font=MONO, size=10, color=INK_MUTE,
                 align=PP_ALIGN.CENTER)
        cx = cx + render_w + Inches(0.35)


def slide_current(prs):
    s = slide_chrome(prs)
    add_text(s, Inches(0.9), Inches(0.5), Inches(8), Inches(0.4),
             "0X · CURRENT MARK (FOR REFERENCE)",
             font=MONO, size=11, color=C_PRIM)
    add_text(s, Inches(0.9), Inches(0.9), Inches(14), Inches(1.2),
             "Syringe on target.",
             font=DISPLAY, size=42, bold=True, color=INK)
    add_text(s, Inches(0.9), Inches(2.4), Inches(8), Inches(3),
             "The current logomark from the github.com/needlegps org "
             "avatar. Complex (syringe at angle, partial target rings, "
             "crosshair ticks, two colors) and reads as detailed at hero "
             "size, but blurs at favicon (32px) scale.\n\n"
             "Brand decision: keep, or replace with one of the simpler "
             "options A–F.",
             font=DISPLAY, size=15, color=INK_SOFT)

    sizes = [(256, Inches(2.4)), (96, Inches(1.0)), (32, Inches(0.42))]
    cx = Inches(9.5); cy = Inches(2.4)
    for px, render_w in sizes:
        bio = io.BytesIO(CURRENT_MARK.read_bytes())
        s.shapes.add_picture(bio, cx, cy, height=render_w, width=render_w)
        add_text(s, cx, cy + render_w + Inches(0.05), render_w, Inches(0.3),
                 f"{px}px", font=MONO, size=10, color=INK_MUTE,
                 align=PP_ALIGN.CENTER)
        cx = cx + render_w + Inches(0.35)


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slide_cover(prs)
    slide_compare(prs)
    slide_current(prs)

    options = [
        ("A", "Bare reticle.",
         LOGO_DIR / "option-A-bare-reticle.svg",
         "Three concentric circles with a center dot. Most stripped-down "
         "version of the current mark. Reads as a target instantly and "
         "stays crisp at 16-32 px favicon scale. Risk: indistinguishable "
         "from a generic 'aim' icon if used alone without the NeedleGPS "
         "wordmark."),
        ("B", "Crosshair.",
         LOGO_DIR / "option-B-crosshair.svg",
         "Ring with four cardinal ticks crossing the circumference. Sharper "
         "than option A; classic targeting reticle vocabulary. The four "
         "ticks add legibility at small sizes by breaking the silhouette."),
        ("C", "Needle arrow.",
         LOGO_DIR / "option-C-needle-arrow.svg",
         "A bold diagonal needle line ending at the target. Conveys 'guided "
         "to target' with a single gesture. More narrative than A or B; "
         "slightly more complex but still single-color and single-line."),
        ("D", "Bubble level.",
         LOGO_DIR / "option-D-bubble-level.svg",
         "Outer ring with an offset filled inner circle representing the "
         "bubble's position. The device's actual physical metaphor. Most "
         "differentiated from generic targeting icons; reads as 'level' or "
         "'bullseye' depending on viewing context."),
        ("E", "Trajectory.",
         LOGO_DIR / "option-E-trajectory.svg",
         "Dashed line from upper-left to a small target near center-bottom. "
         "Smallest visual footprint of the set. Asymmetric composition can "
         "feel 'designed for a logo' but may read as cluttered at favicon "
         "scale."),
        ("F", "N monogram.",
         LOGO_DIR / "option-F-monogram.svg",
         "Geometric ‘N’ letterform with a small cyan dot at its "
         "center. The only option that's a letterform first, mark second. "
         "Strong wordmark integration; loses the targeting metaphor unless "
         "paired with the full NeedleGPS wordmark beside it."),
    ]
    for i, (label, name, svg, rationale) in enumerate(options):
        slide_for_option(prs, i, label, name, svg, rationale)

    out = WEBSITE.parent / "corp docs" / "NeedleGPS_Logo_Options.pptx"
    prs.save(str(out))
    print(f"wrote {out}  ({out.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
